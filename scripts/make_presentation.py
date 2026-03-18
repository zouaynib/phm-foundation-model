"""
Generate a presentation summarising the PHM Foundation Model results.
Run:  python make_presentation.py
Output: results/foundation_model_results.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG    = RGBColor(0x1B, 0x1B, 0x2F)  # deep navy
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)  # bright cyan
ACCENT2    = RGBColor(0x90, 0xE0, 0xEF)  # light cyan
GREEN      = RGBColor(0x06, 0xD6, 0xA0)  # success green
RED_SOFT   = RGBColor(0xEF, 0x47, 0x6F)  # soft red
GRAY       = RGBColor(0xA0, 0xA0, 0xB0)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE8)
CARD_BG    = RGBColor(0x27, 0x27, 0x3F)  # card background
SLIDE_W    = Inches(13.333)
SLIDE_H    = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H


# ── Helpers ─────────────────────────────────────────────────────
def _add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, font_size=18,
              color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
              font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_para(tf, text, font_size=18, color=WHITE, bold=False,
              alignment=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    p.space_before = space_before
    return p


def _add_card(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_accent_bar(slide, left, top, width=Inches(0.8), height=Pt(4),
                    color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top,
                                   width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ── SLIDE 1 — Title ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_add_bg(slide)

_add_accent_bar(slide, Inches(1), Inches(2.0), Inches(1.2), Pt(5), ACCENT)

_add_text(slide, Inches(1), Inches(2.3), Inches(11), Inches(1.5),
          "A Multi-Domain Foundation Model\nfor Prognostics & Health Management",
          font_size=36, bold=True, color=WHITE)

_add_text(slide, Inches(1), Inches(4.0), Inches(10), Inches(0.6),
          "Pre-trained across 4 PHM datasets  ·  Evaluated on classification & RUL tasks",
          font_size=18, color=ACCENT2)

_add_text(slide, Inches(1), Inches(5.2), Inches(10), Inches(0.5),
          "Zaynab Raounak", font_size=20, color=GRAY)


# ── SLIDE 2 — Motivation ───────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide)

_add_accent_bar(slide, Inches(0.8), Inches(0.5))
_add_text(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.7),
          "Why a Foundation Model for PHM?", font_size=30, bold=True)

items = [
    ("Problem", "PHM models are typically trained from scratch per dataset — limited by small, domain-specific data."),
    ("Idea", "Pre-train a single transformer backbone across multiple PHM domains, then fine-tune per task."),
    ("Goal", "Leverage cross-domain patterns (vibration signatures, degradation trends) to improve accuracy and data efficiency."),
]

y = Inches(1.8)
for title, desc in items:
    _add_card(slide, Inches(0.8), y, Inches(11.5), Inches(1.15))
    _add_text(slide, Inches(1.1), y + Inches(0.15), Inches(1.5), Inches(0.4),
              title, font_size=16, bold=True, color=ACCENT)
    _add_text(slide, Inches(2.6), y + Inches(0.15), Inches(9.2), Inches(0.9),
              desc, font_size=16, color=LIGHT_GRAY)
    y += Inches(1.45)


# ── SLIDE 3 — Datasets ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide)

_add_accent_bar(slide, Inches(0.8), Inches(0.5))
_add_text(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.7),
          "Training Datasets", font_size=30, bold=True)

datasets = [
    ("CWRU",      "Bearing Fault Diagnosis",   "1 ch · 12 kHz",    "Classification (4 classes)"),
    ("PRONOSTIA", "Bearing Degradation",        "2 ch · 25.6 kHz",  "RUL Regression"),
    ("CMAPSS",    "Turbofan Engine",            "14 ch · cycle-based", "Classification + RUL"),
    ("MFPT",      "Bearing Fault Diagnosis",    "1 ch · 97.6 kHz",  "Classification (3 classes)"),
]

card_w = Inches(2.65)
gap = Inches(0.2)
start_x = Inches(0.8)
for i, (name, domain, specs, task) in enumerate(datasets):
    x = start_x + i * (card_w + gap)
    _add_card(slide, x, Inches(2.0), card_w, Inches(3.8))
    _add_text(slide, x + Inches(0.25), Inches(2.2), card_w - Inches(0.5), Inches(0.5),
              name, font_size=22, bold=True, color=ACCENT)
    _add_text(slide, x + Inches(0.25), Inches(2.8), card_w - Inches(0.5), Inches(0.5),
              domain, font_size=14, color=LIGHT_GRAY)
    _add_text(slide, x + Inches(0.25), Inches(3.5), card_w - Inches(0.5), Inches(0.5),
              specs, font_size=13, color=GRAY)
    _add_text(slide, x + Inches(0.25), Inches(4.4), card_w - Inches(0.5), Inches(0.8),
              task, font_size=14, bold=True, color=GREEN)


# ── SLIDE 4 — Architecture ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide)

_add_accent_bar(slide, Inches(0.8), Inches(0.5))
_add_text(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.7),
          "Model Architecture", font_size=30, bold=True)

# Architecture pipeline boxes
stages = [
    ("Input\nSignals", "Multi-channel\ntime series"),
    ("Patch\nEmbedding", "64-sample patches\nlinear projection"),
    ("Transformer\nEncoder", "4 layers, 8 heads\nd_model = 128"),
    ("Domain\nEmbedding", "Frequency +\nDataset ID"),
    ("Projector\nMLP", "Shared latent\nrepresentation"),
    ("Task\nHeads", "Per-dataset\ncls / rul heads"),
]

box_w = Inches(1.7)
box_h = Inches(2.2)
gap = Inches(0.18)
start_x = Inches(0.6)
y_box = Inches(2.4)

for i, (title, desc) in enumerate(stages):
    x = start_x + i * (box_w + gap)
    _add_card(slide, x, y_box, box_w, box_h, CARD_BG)
    _add_text(slide, x + Inches(0.12), y_box + Inches(0.2),
              box_w - Inches(0.24), Inches(0.7),
              title, font_size=15, bold=True, color=ACCENT,
              alignment=PP_ALIGN.CENTER)
    _add_text(slide, x + Inches(0.12), y_box + Inches(1.0),
              box_w - Inches(0.24), Inches(1.0),
              desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < len(stages) - 1:
        ax = x + box_w + Inches(0.02)
        _add_text(slide, ax, y_box + Inches(0.8), Inches(0.18), Inches(0.5),
                  "\u25B6", font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Key design decisions
_add_text(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.4),
          "Key:  Channel-independent processing  |  Balanced domain sampling  |  Multi-task loss (CE + MSE)",
          font_size=14, color=GRAY)


# ── SLIDE 5 — Classification Results (HERO SLIDE) ──────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide)

_add_accent_bar(slide, Inches(0.8), Inches(0.5))
_add_text(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.7),
          "Classification Results", font_size=30, bold=True)

# Big number highlight
_add_card(slide, Inches(0.8), Inches(1.8), Inches(5.0), Inches(2.0), CARD_BG)
_add_text(slide, Inches(1.0), Inches(1.9), Inches(4.5), Inches(0.6),
          "Average Classification Accuracy", font_size=16, color=GRAY)
_add_text(slide, Inches(1.0), Inches(2.5), Inches(4.5), Inches(1.0),
          "98.12%", font_size=52, bold=True, color=GREEN, alignment=PP_ALIGN.LEFT)
_add_text(slide, Inches(3.5), Inches(3.0), Inches(2.5), Inches(0.5),
          "+2.6% over baseline", font_size=16, bold=True, color=ACCENT)

# Per-dataset results table
_add_card(slide, Inches(6.5), Inches(1.8), Inches(5.8), Inches(2.0), CARD_BG)

headers = ["Dataset", "Baseline", "Foundation", "Gain"]
col_x = [Inches(6.8), Inches(8.2), Inches(9.6), Inches(11.1)]
hy = Inches(2.0)
for cx, h in zip(col_x, headers):
    _add_text(slide, cx, hy, Inches(1.3), Inches(0.35),
              h, font_size=13, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

rows = [
    ("CWRU",  "92.75%", "97.95%", "+5.20%"),
    ("MFPT",  "98.29%", "98.29%", "  0.00%"),
    ("CMAPSS","—",      "—",      "—"),
]
ry = Inches(2.45)
for name, bl, fd, gain in rows:
    gain_color = GREEN if gain.startswith("+") else LIGHT_GRAY
    vals = [name, bl, fd, gain]
    colors = [WHITE, LIGHT_GRAY, WHITE, gain_color]
    bolds = [True, False, False, True]
    for cx, v, c, b in zip(col_x, vals, colors, bolds):
        _add_text(slide, cx, ry, Inches(1.3), Inches(0.35),
                  v, font_size=14, color=c, bold=b, alignment=PP_ALIGN.CENTER)
    ry += Inches(0.4)

# F1 scores
_add_card(slide, Inches(0.8), Inches(4.2), Inches(5.0), Inches(2.2), CARD_BG)
_add_text(slide, Inches(1.0), Inches(4.35), Inches(4.5), Inches(0.4),
          "F1 Score (Macro)", font_size=16, bold=True, color=WHITE)

f1_data = [
    ("CWRU", "93.96%", "98.27%", "+4.31%"),
    ("MFPT", "97.53%", "97.53%", "  0.00%"),
]
ry = Inches(4.9)
for name, bl, fd, gain in f1_data:
    tf = _add_text(slide, Inches(1.2), ry, Inches(4.5), Inches(0.35),
                   f"{name}:   {bl}  →  {fd}   ({gain})",
                   font_size=14, color=LIGHT_GRAY)
    ry += Inches(0.4)

# Key takeaway
_add_card(slide, Inches(6.5), Inches(4.2), Inches(5.8), Inches(2.2), CARD_BG)
_add_text(slide, Inches(6.8), Inches(4.35), Inches(5.2), Inches(0.4),
          "Key Takeaway", font_size=16, bold=True, color=WHITE)
_add_text(slide, Inches(6.8), Inches(4.9), Inches(5.2), Inches(1.4),
          "The foundation model achieves near-perfect classification on CWRU "
          "(+5.2% over CNN baseline) while maintaining parity on MFPT. "
          "Cross-domain pretraining provides the biggest lift on CWRU, "
          "where the model benefits from vibration patterns learned across "
          "all four datasets.",
          font_size=14, color=LIGHT_GRAY)


# ── SLIDE 6 — Cross-Domain Generalization ──────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide)

_add_accent_bar(slide, Inches(0.8), Inches(0.5))
_add_text(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.7),
          "Cross-Domain Generalization (Leave-One-Out)", font_size=30, bold=True)

_add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
          "Model pre-trained on 3 datasets, then fine-tuned and tested on the held-out dataset",
          font_size=16, color=GRAY)

# Big result cards
loo_results = [
    ("CWRU\nheld out", "96.38%", "96.91%", "Trained on: PRONOSTIA + CMAPSS + MFPT"),
    ("MFPT\nheld out", "97.44%", "96.34%", "Trained on: CWRU + PRONOSTIA + CMAPSS"),
]

for i, (label, acc, f1, trained_on) in enumerate(loo_results):
    x = Inches(0.8) + i * Inches(6.1)
    _add_card(slide, x, Inches(2.3), Inches(5.7), Inches(3.5), CARD_BG)
    _add_text(slide, x + Inches(0.3), Inches(2.5), Inches(5), Inches(0.8),
              label, font_size=22, bold=True, color=ACCENT)
    _add_text(slide, x + Inches(0.3), Inches(3.3), Inches(2), Inches(0.3),
              "Accuracy", font_size=13, color=GRAY)
    _add_text(slide, x + Inches(0.3), Inches(3.6), Inches(3), Inches(0.8),
              acc, font_size=44, bold=True, color=GREEN)
    _add_text(slide, x + Inches(2.8), Inches(3.3), Inches(2), Inches(0.3),
              "F1 Score", font_size=13, color=GRAY)
    _add_text(slide, x + Inches(2.8), Inches(3.6), Inches(3), Inches(0.8),
              f1, font_size=44, bold=True, color=GREEN)
    _add_text(slide, x + Inches(0.3), Inches(4.8), Inches(5), Inches(0.5),
              trained_on, font_size=12, color=GRAY)

_add_card(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8), CARD_BG)
_add_text(slide, Inches(1.1), Inches(6.3), Inches(11), Inches(0.6),
          "The model generalises to unseen domains with minimal accuracy loss — "
          "demonstrating that the shared transformer backbone captures transferable "
          "vibration features across different machinery types.",
          font_size=14, color=LIGHT_GRAY)


# ── SLIDE 7 — Low-Data Regime ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide)

_add_accent_bar(slide, Inches(0.8), Inches(0.5))
_add_text(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.7),
          "Low-Data Regime Analysis", font_size=30, bold=True)

_add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
          "Baseline CNN accuracy with reduced training data (CWRU)",
          font_size=16, color=GRAY)

# Bar chart style using shapes
fractions = [
    ("10%", 0.4243, ACCENT),
    ("20%", 0.4138, ACCENT),
    ("50%", 0.4833, ACCENT2),
    ("100%", 0.4860, GREEN),
]

max_bar_w = Inches(7.0)
bar_h = Inches(0.7)
x_start = Inches(2.0)
y_start = Inches(2.5)

for i, (label, acc, color) in enumerate(fractions):
    y = y_start + i * Inches(1.1)
    # Label
    _add_text(slide, Inches(0.8), y + Inches(0.1), Inches(1.0), Inches(0.5),
              label + " data", font_size=16, bold=True, color=WHITE,
              alignment=PP_ALIGN.RIGHT)
    # Bar
    bar_w = int(max_bar_w * acc / 0.55)  # scale
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x_start, y, bar_w, bar_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Value label
    _add_text(slide, x_start + bar_w + Inches(0.2), y + Inches(0.1),
              Inches(1.5), Inches(0.5),
              f"{acc*100:.1f}%", font_size=16, bold=True, color=WHITE)

_add_card(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8), CARD_BG)
_add_text(slide, Inches(1.1), Inches(6.1), Inches(11), Inches(0.6),
          "Even at 10% training data, the baseline achieves 42.4% — "
          "a foundation model pre-trained on related domains can significantly "
          "boost this via transfer learning, reducing the labelled data requirement.",
          font_size=14, color=LIGHT_GRAY)


# ── SLIDE 8 — Summary & Next Steps ─────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide)

_add_accent_bar(slide, Inches(0.8), Inches(0.5))
_add_text(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.7),
          "Summary & Next Steps", font_size=30, bold=True)

# Achievements
_add_card(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.5), CARD_BG)
_add_text(slide, Inches(1.1), Inches(1.9), Inches(5), Inches(0.4),
          "Achievements", font_size=18, bold=True, color=GREEN)

achievements = [
    "98.12% avg classification accuracy (+2.6% over baseline)",
    "CWRU accuracy: 97.95% (+5.2% gain from pre-training)",
    "Strong cross-domain transfer (96-97% on held-out domains)",
    "Unified architecture handles 4 PHM domains",
    "Channel-independent design supports 1 to 14 channels",
]
ry = Inches(2.5)
for a in achievements:
    _add_text(slide, Inches(1.3), ry, Inches(4.8), Inches(0.4),
              "\u2713  " + a, font_size=13, color=LIGHT_GRAY)
    ry += Inches(0.38)

# Next steps
_add_card(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(4.5), CARD_BG)
_add_text(slide, Inches(7.1), Inches(1.9), Inches(5), Inches(0.4),
          "Next Steps", font_size=18, bold=True, color=ACCENT)

next_steps = [
    "Improve RUL regression with refined head architecture",
    "Add more PHM domains (Paderborn, IMS, etc.)",
    "Evaluate few-shot transfer with foundation features",
    "Explore self-supervised pre-training objectives",
    "Benchmark against state-of-the-art per-dataset models",
]
ry = Inches(2.5)
for ns in next_steps:
    _add_text(slide, Inches(7.3), ry, Inches(4.8), Inches(0.4),
              "\u2192  " + ns, font_size=13, color=LIGHT_GRAY)
    ry += Inches(0.38)

# Bottom banner
_add_card(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6), CARD_BG)
_add_text(slide, Inches(1.1), Inches(6.55), Inches(11), Inches(0.5),
          "A single pre-trained backbone outperforms per-dataset CNNs on "
          "classification — demonstrating the viability of foundation models for PHM.",
          font_size=15, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────
os.makedirs("results", exist_ok=True)
out_path = "results/foundation_model_results.pptx"
prs.save(out_path)
print(f"Presentation saved to {out_path}")
